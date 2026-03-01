import json, sys, io, requests, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.chart import XL_TICK_MARK
import cairosvg

if len(sys.argv) < 3: 
    print("Usage: python generate_ppt.py <input.json> <output.pptx> [template.pptx]")
    sys.exit(1)

input_json = sys.argv[1]
output_pptx = sys.argv[2]
script_dir = os.path.dirname(os.path.abspath(__file__))

# Use provided template, or the embedded base_template.pptx, or default fallback
if len(sys.argv) >= 4:
    template_path = sys.argv[3]
else:
    template_path = os.path.join(script_dir, 'base_template.pptx')

if os.path.exists(template_path):
    prs = Presentation(template_path)
    # Clear existing slides from the template
    for i in range(len(prs.slides) - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]
else:
    prs = Presentation() # Fallback to blank
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)

with open(input_json, 'r', encoding='utf-8') as f: 
    slides_data = json.load(f)

# Corporate Design System Colors
BG = RGBColor(255, 255, 255)
TITLE_COLOR = RGBColor(0, 82, 155)      # Corporate Deep Blue
ACCENT = RGBColor(0, 153, 204)          # Light Blue
TEXT = RGBColor(60, 60, 60)
BOX_BG = RGBColor(245, 247, 250)
ICON_COLOR_HEX = "00529B"

PALETTE = [
    RGBColor(0, 82, 155),    # Deep Blue
    RGBColor(0, 153, 204),   # Light Blue
    RGBColor(255, 153, 0),   # Orange
    RGBColor(0, 153, 102),   # Teal
    RGBColor(153, 0, 51),    # Dark Red
    RGBColor(102, 102, 102), # Dark Gray
    RGBColor(153, 102, 204), # Purple
]

def download_icon(icon_name):
    try:
        if ":" not in icon_name:
            icon_name = f"lucide:{icon_name}"
        prefix, name = icon_name.split(":")
        url = f"https://api.iconify.design/{prefix}/{name}.svg?color=%23{ICON_COLOR_HEX}&width=128"
        svg_data = requests.get(url, timeout=5).content
        png_data = cairosvg.svg2png(bytestring=svg_data)
        return io.BytesIO(png_data)
    except Exception:
        return None

def add_bullet(tf, text):
    p = tf.paragraphs[0] if not tf.text else tf.add_paragraph()
    p.level, p.line_spacing, p.space_after = 0, 1.3, Pt(12)
    colon = text.find("：") if "：" in text else text.find(": ")
    if colon != -1:
        r1, r2 = p.add_run(), p.add_run()
        r1.text, r1.font.name, r1.font.size, r1.font.bold, r1.font.color.rgb = text[:colon+1], 'Microsoft YaHei', Pt(15), True, TEXT
        r2.text, r2.font.name, r2.font.size, r2.font.color.rgb = text[colon+1:], 'Microsoft YaHei', Pt(15), TEXT
    else:
        r = p.add_run()
        r.text, r.font.name, r.font.size, r.font.color.rgb = text, 'Microsoft YaHei', Pt(15), TEXT

for data in slides_data:
    layout = data.get("layout", "two-column")
    try:
        slide_layout = prs.slide_layouts[6] # Often a blank layout
    except:
        slide_layout = prs.slide_layouts[0]
        
    slide = prs.slides.add_slide(slide_layout)
    # Ensure background is white if not defined by master
    try:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BG
    except: pass
    
    if layout == "cover":
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(2))
        tf = txBox.text_frame; tf.word_wrap = True; tf.auto_size = 2
        p = tf.paragraphs[0]
        p.text, p.font.name, p.font.size, p.font.bold, p.font.color.rgb, p.alignment = data.get("title", ""), 'Microsoft YaHei', Pt(40), True, TITLE_COLOR, PP_ALIGN.CENTER
        if "subtitle" in data:
            p2 = tf.add_paragraph()
            p2.text, p2.font.name, p2.font.size, p2.font.color.rgb, p2.alignment, p2.space_before = data.get("subtitle", ""), 'Microsoft YaHei', Pt(24), ACCENT, PP_ALIGN.CENTER, Pt(20)
        continue

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
    tf_title = title_box.text_frame; tf_title.word_wrap = True; tf_title.auto_size = 2
    p = tf_title.paragraphs[0]
    p.text, p.font.name, p.font.size, p.font.bold, p.font.color.rgb = data.get("action_title", ""), 'Microsoft YaHei', Pt(20), True, TITLE_COLOR
    slide.shapes.add_connector(1, Inches(0.5), Inches(1.3), Inches(12.83), Inches(1.3)).line.color.rgb = ACCENT
    
    content_bottom = 6.8
    if "takeaway" in data:
        content_bottom = 5.9
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.1), Inches(12.33), Inches(0.9))
        shape.fill.solid(); shape.fill.fore_color.rgb = BOX_BG
        shape.line.color.rgb, shape.line.width = ACCENT, Pt(1.5)
        tf_ta = shape.text_frame; tf_ta.word_wrap = True; tf_ta.auto_size = 2
        p_ta = tf_ta.paragraphs[0]
        p_ta.text, p_ta.font.name, p_ta.font.size, p_ta.font.bold, p_ta.font.color.rgb = data["takeaway"], 'Microsoft YaHei', Pt(16), True, TITLE_COLOR

    if layout in ["two-column", "three-column"]:
        cols = data.get("columns", [])
        n_cols = len(cols)
        if n_cols > 0:
            spacing, total_width = 0.4, 12.33
            col_width = (total_width - (n_cols - 1) * spacing) / n_cols
            for i, col in enumerate(cols):
                x, y, h = 0.5 + i * (col_width + spacing), 1.5, content_bottom - 1.5
                icon_name = col.get("icon")
                icon_offset = 0
                if icon_name:
                    img_stream = download_icon(icon_name)
                    if img_stream:
                        slide.shapes.add_picture(img_stream, Inches(x + 0.1), Inches(y + 0.05), width=Inches(0.4), height=Inches(0.4))
                        icon_offset = 0.5
                if "title" in col:
                    ch_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(col_width), Inches(0.5))
                    ch_box.fill.solid(); ch_box.fill.fore_color.rgb = BOX_BG
                    ch_box.line.color.rgb = TITLE_COLOR; ch_box.line.width = Pt(1)
                    tt_box = slide.shapes.add_textbox(Inches(x + icon_offset), Inches(y), Inches(col_width - icon_offset), Inches(0.5))
                    tf_ch = tt_box.text_frame; tf_ch.vertical_anchor = 3
                    p_ch = tf_ch.paragraphs[0]
                    p_ch.text, p_ch.font.name, p_ch.font.size, p_ch.font.bold, p_ch.font.color.rgb, p_ch.alignment = col["title"], 'Microsoft YaHei', Pt(16), True, TITLE_COLOR, PP_ALIGN.LEFT
                    y += 0.6; h -= 0.6
                body_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(col_width), Inches(h))
                tf_body = body_box.text_frame; tf_body.word_wrap = True; tf_body.auto_size = 2
                for bullet in col.get("bullets", []): add_bullet(tf_body, "• " + bullet)
    
    elif layout == "image-text":
        img_path = data.get("image_path")
        if img_path and os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), width=Inches(6.0), height=Inches(4.2))
            body_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(6.0), Inches(4.2))
            tf_body = body_box.text_frame; tf_body.word_wrap = True; tf_body.auto_size = 2
            for content_line in data.get("content", []): add_bullet(tf_body, "• " + content_line)
        else:
            body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.33), Inches(4.2))
            tf_body = body_box.text_frame; tf_body.word_wrap = True; tf_body.auto_size = 2
            for content_line in data.get("content", []): add_bullet(tf_body, "• " + content_line)

    elif layout == "timeline":
        steps = data.get("steps", [])
        n_steps = len(steps)
        if n_steps > 0:
            spacing, start_x, total_avail_width = 0.3, 0.5, 12.33
            width = (total_avail_width - (n_steps - 1) * spacing) / n_steps
            for i, step in enumerate(steps):
                x, shape_y, shape_h = start_x + i * (width + spacing), 2.5, 0.8
                shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(shape_y), Inches(width), Inches(shape_h))
                shape.fill.solid(); shape.fill.fore_color.rgb = TITLE_COLOR if i == 2 else RGBColor(230, 230, 230)
                shape.line.fill.background()
                tf = shape.text_frame; tf.word_wrap = True
                p = tf.paragraphs[0]; r = p.add_run()
                r.text, r.font.size, r.font.bold, r.font.color.rgb = step.get("title", ""), Pt(14), True, BG if i == 2 else TEXT
                p.alignment = PP_ALIGN.CENTER
                icon_name = step.get("icon")
                icon_y = shape_y + shape_h + 0.3
                if icon_name:
                    img_stream = download_icon(icon_name)
                    if img_stream:
                        slide.shapes.add_picture(img_stream, Inches(x + (width-0.6)/2), Inches(icon_y), width=Inches(0.6), height=Inches(0.6))
                        icon_y += 0.8
                tb = slide.shapes.add_textbox(Inches(x), Inches(icon_y), Inches(width), Inches(content_bottom - icon_y))
                tb.text_frame.word_wrap = True; tb.text_frame.auto_size = 2
                p2 = tb.text_frame.paragraphs[0]; r2 = p2.add_run()
                r2.text, r2.font.size, r2.font.color.rgb = step.get("desc", ""), Pt(13), TEXT

    elif layout == "matrix":
        quads = data.get("quadrants", [])
        coords = [(1.5, 1.6), (7, 1.6), (1.5, 3.8), (7, 3.8)]
        for i, quad in enumerate(quads):
            if i >= len(coords): break
            x, y = coords[i]
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.0), Inches(2.0))
            shape.fill.solid(); shape.fill.fore_color.rgb = BOX_BG; shape.line.color.rgb = ACCENT
            tf = shape.text_frame; tf.word_wrap = True; tf.auto_size = 2
            tf.margin_left = Inches(0.8)
            p = tf.paragraphs[0]
            p.text, p.font.bold, p.font.size, p.font.color.rgb = quad.get("title", ""), True, Pt(16), TITLE_COLOR
            p2 = tf.add_paragraph(); r2 = p2.add_run()
            r2.text = "\\n" + quad.get("desc", "")
            r2.font.size = Pt(14)
            r2.font.color.rgb = TEXT
            icon_name = quad.get("icon")
            if icon_name:
                img_stream = download_icon(icon_name)
                if img_stream: slide.shapes.add_picture(img_stream, Inches(x + 0.15), Inches(y + 0.2), width=Inches(0.5), height=Inches(0.5))

    elif layout == "native-chart":
        c_type = data.get("chart_type", "column_clustered")
        categories = data.get("categories", [])
        series_list = data.get("series", [])
        
        max_cat_len = max([len(str(c)) for c in categories]) if categories else 0
        num_categories = len(categories)
        num_series = len(series_list)
        total_data_points = num_categories * num_series
        
        if (max_cat_len > 6 and num_categories > 4) or num_categories > 8:
            c_type = "bar_clustered"
            
        chart_data = ChartData() if c_type == "pie" else CategoryChartData()
        chart_data.categories = categories
        
        for series in series_list:
            chart_data.add_series(series.get("name", "Value"), series.get("values", []))
            
        x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(12.33), Inches(4.5)
        
        type_enum = XL_CHART_TYPE.COLUMN_CLUSTERED
        if c_type == "bar_clustered": type_enum = XL_CHART_TYPE.BAR_CLUSTERED
        if c_type == "pie": type_enum = XL_CHART_TYPE.PIE
        
        def style_chart(chart, c_type, num_series, total_data_points, max_cat_len):
            # AESTHETICS: Thick bars (Gap Width = 100)
            try:
                if chart.chart_type in [XL_CHART_TYPE.COLUMN_CLUSTERED, XL_CHART_TYPE.BAR_CLUSTERED]:
                    chart.plots[0].gap_width = 100
            except: pass
            
            # AESTHETICS: Palette Application
            if c_type == "pie":
                chart.plots[0].vary_by_categories = True
                try:
                    for idx, point in enumerate(chart.series[0].points):
                        fill = point.format.fill
                        fill.solid()
                        fill.fore_color.rgb = PALETTE[idx % len(PALETTE)]
                except: pass
            else:
                for idx_s, series in enumerate(chart.series):
                    fill = series.format.fill
                    fill.solid()
                    fill.fore_color.rgb = PALETTE[idx_s % len(PALETTE)]
                    series.format.line.fill.background() # No outlines
                    
            # AESTHETICS: Legend positioning
            chart.has_legend = True if num_series > 1 or c_type == "pie" else False
            if chart.has_legend:
                chart.legend.position = XL_LEGEND_POSITION.RIGHT if num_series > 4 else XL_LEGEND_POSITION.BOTTOM
                chart.legend.font.size = Pt(10)
                chart.legend.font.name = 'Microsoft YaHei'
                chart.legend.include_in_layout = False
                
            # AESTHETICS: Data Labels
            chart.plots[0].has_data_labels = True
            data_labels = chart.plots[0].data_labels
            data_labels.font.name = 'Microsoft YaHei'
            
            # Data Type detection: score vs percentage
            data_type = data.get("data_type", "percentage")
            if data_type == "score":
                data_labels.number_format = '0.00'
            elif c_type == "pie":
                data_labels.number_format = '0%'
            else:
                data_labels.number_format = '0"%"'
                
            data_labels.number_format_is_linked = False
            
            try:
                data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
            except: pass
            
            # Font size scaling
            if total_data_points > 30:
                data_labels.font.size = Pt(10)
            elif total_data_points > 15:
                data_labels.font.size = Pt(12)
            else:
                data_labels.font.size = Pt(14)
            
            if total_data_points > 40:
                chart.plots[0].has_data_labels = False 
                
            if c_type == "pie":
                data_labels.show_percentage = True
                data_labels.show_value = False
                
            # AESTHETICS: Clean Axes (No lines, specific font sizes)
            try:
                cat_axis = chart.category_axis
                cat_axis.format.line.fill.background()
                cat_axis.tick_labels.font.name = 'Microsoft YaHei'
                if max_cat_len > 12:
                    cat_axis.tick_labels.font.size = Pt(8)
                elif max_cat_len > 5:
                    cat_axis.tick_labels.font.size = Pt(9)
                else:
                    cat_axis.tick_labels.font.size = Pt(10)
                cat_axis.tick_labels.font.color.rgb = RGBColor(80, 80, 80)
                cat_axis.tick_label_spacing = 1
            except ValueError:
                pass
                
            try:
                val_axis = chart.value_axis
                val_axis.format.line.fill.background()
                val_axis.has_major_gridlines = False # CRITICAL: No gridlines per user requirement
                val_axis.tick_labels.font.name = 'Microsoft YaHei'
                val_axis.tick_labels.font.size = Pt(9)
                val_axis.tick_labels.font.color.rgb = RGBColor(120, 120, 120)
            except ValueError:
                pass

        # SPLIT LOGIC FOR DENSE BAR CHARTS
        if type_enum == XL_CHART_TYPE.BAR_CLUSTERED and len(categories) >= 10 and num_series == 1:
            mid = len(categories) // 2
            cd1 = CategoryChartData()
            cd1.categories = categories[:mid]
            cd1.add_series(series_list[0].get("name", "Value"), series_list[0].get("values", [])[:mid])
            cd2 = CategoryChartData()
            cd2.categories = categories[mid:]
            cd2.add_series(series_list[0].get("name", "Value"), series_list[0].get("values", [])[mid:])
            
            chart1 = slide.shapes.add_chart(type_enum, Inches(0.2), y, Inches(6.0), cy, cd1).chart
            style_chart(chart1, c_type, num_series, total_data_points//2, max_cat_len)
            
            chart2 = slide.shapes.add_chart(type_enum, Inches(6.5), y, Inches(6.5), cy, cd2).chart
            style_chart(chart2, c_type, num_series, total_data_points//2, max_cat_len)
            
        elif type_enum == XL_CHART_TYPE.BAR_CLUSTERED and max_cat_len > 10:
            cx_adj = Inches(10.5)
            x_adj = Inches(2.0)
            chart = slide.shapes.add_chart(type_enum, x_adj, y, cx_adj, cy, chart_data).chart
            style_chart(chart, c_type, num_series, total_data_points, max_cat_len)
        else:
            chart = slide.shapes.add_chart(type_enum, x, y, cx, cy, chart_data).chart
            style_chart(chart, c_type, num_series, total_data_points, max_cat_len)

prs.save(output_pptx)
